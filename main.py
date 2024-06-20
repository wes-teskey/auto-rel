import csv
import random
import numpy as np
from utils import path_to_save_file
import pandas as pd
import matplotlib.pyplot as plt
from crewai import Agent, Task, Crew, Process
from crewai.project import CrewBase, agent, crew, task
from dotenv import load_dotenv
from langchain.chat_models import ChatOpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

load_dotenv()

@CrewBase
class RelReportCrew():
    agents_config = 'rel_agents.yaml'
    tasks_config = 'rel_tasks.yaml'
    
    def __init__(self, test_1_file_name, test_2_file_name, 
                 test_3_file_name, report_summary_file_name, llm):
        self.test_1_file_name = test_1_file_name
        self.test_2_file_name = test_2_file_name
        self.test_3_file_name = test_3_file_name
        self.report_summary_file_name = report_summary_file_name
        self.llm = llm
  
    @agent
    def rel_engineer(self) -> Agent:
        return Agent(
          config = self.agents_config['rel_engineer'],
          llm = self.llm,
          )
        
    @agent
    def rel_manager(self) -> Agent:
        return Agent(
          config = self.agents_config['rel_manager'],
          llm = self.llm,
          )
    
    @task
    def test_1_summary(self) -> Task:
        return Task(
            config = self.tasks_config['test_1_summary'],
            agent = self.rel_engineer(),
            context=[],
            output_file = self.test_1_file_name,
            )
    
    @task
    def test_2_summary(self) -> Task:
        return Task(
            config = self.tasks_config['test_2_summary'],
            agent = self.rel_engineer(),
            context=[],
            output_file = self.test_2_file_name,
            )
    
    @task
    def test_3_summary(self) -> Task:
        return Task(
            config = self.tasks_config['test_3_summary'],
            agent = self.rel_engineer(),
            context=[],
            output_file = self.test_3_file_name,
            )
        
    @task
    def report_summary(self) -> Task:
        return Task(
            config = self.tasks_config['report_summary'],
            agent = self.rel_manager(),
            context=[self.test_1_summary(), self.test_2_summary(), self.test_3_summary()],
            output_file = self.report_summary_file_name,
            )
    
    @crew
    def crew(self) -> Crew:
        return Crew(
            agents = [self.rel_engineer(), self.rel_manager()],
            tasks = [self.test_1_summary(), self.test_2_summary(),
                     self.test_3_summary(), self.report_summary()],
            process = Process.sequential,
            verbose = 2,
            )

# Generate unique die IDs randomly assigned from 10,000 to 99,999
def generate_unique_ids(n, start, end):
    return random.sample(range(start, end), n)

# Generate Idss values from two normal distributions
def generate_idss_values(mean1, sigma1, mean2, sigma2, size1, size2):
    idss_values = np.concatenate([
        np.random.normal(mean1, sigma1, size1),
        np.random.normal(mean2, sigma2, size2)
    ])
    np.random.shuffle(idss_values)
    return np.abs(idss_values)

def create_sample_data(sub_dir, 
                       file_name='reliability_test_data.csv', 
                       column_names=['Die IDs HTOL', 'HTOL Idss (nA)', 
                                     'Die IDs TCT', 'TCT Idss (nA)', 
                                     'Die IDs BHAST', 'BHAST Idss (nA)']):
    # Define the number of rows and number of REL tests
    num_rows = 77
    num_rel_tests = 3

    # Generate die IDs
    all_die_ids = generate_unique_ids(num_rows*num_rel_tests, 10000, 99999)
    die_ids1 = all_die_ids[:num_rows]
    die_ids2 = all_die_ids[num_rows:num_rows*2]
    die_ids3 = all_die_ids[num_rows*2:]

    # Generate Idss values
    idss_values1 = generate_idss_values(0.2, 0.03, 0.5, 0.4, 70, 7)
    idss_values2 = generate_idss_values(0.2, 0.03, 0.5, 0.4, 70, 7)
    idss_values3 = generate_idss_values(0.2, 0.03, 0.5, 0.4, 70, 7)

    # Write to CSV
    file_path = path_to_save_file(sub_dir, file_name, False)
    with open(file_path, 'w', newline='') as file:
        writer = csv.writer(file)
        writer.writerow(column_names)
        for i in range(num_rows):
            writer.writerow([die_ids1[i], idss_values1[i], 
                             die_ids2[i], idss_values2[i], 
                             die_ids3[i], idss_values3[i]])

    print(f"CSV file '{file_name}' has been created.")
    return file_path, column_names

def plot_idss_values(file_path, column_names, sub_dir, Idss_lim):
    # Load the CSV file
    df = pd.read_csv(file_path)

    # Define a function to plot and save histogram
    def plot_histogram(column_name, x_label, y_label, title, file_name, x_max):
        plt.figure()
        plt.hist(df[column_name], bins=20, range=(0, x_max), edgecolor='black')
        plt.axvline(x=Idss_lim, color='red', linestyle='--')
        plt.xlabel(x_label)
        plt.ylabel(y_label)
        plt.title(title)
        plt.xlim(0, x_max)
        plt.grid(True)
        plt.savefig(file_name)
        plt.close()

    #Find names of columns 2, 4, and 6 and create file names
    y_label = 'Number of Samples'
    p1, p2, p3 = column_names[1], column_names[3], column_names[5]
    p1_file = path_to_save_file(sub_dir, p1 + '.png', False)
    p2_file = path_to_save_file(sub_dir, p2 + '.png', False)
    p3_file = path_to_save_file(sub_dir, p3 + '.png', False)

    #Find max value of all plots and Idss_lim for plotting
    x_max = max(df[p1].max(), df[p2].max(), df[p3].max(), Idss_lim)*1.1

    # Plot and save histograms for columns 2, 4, and 6
    plot_histogram(p1, p1.split(' ', 1)[1], y_label, p1.rsplit(' ', 1)[0], p1_file, x_max)
    plot_histogram(p2, p2.split(' ', 1)[1], y_label, p2.rsplit(' ', 1)[0], p2_file, x_max)
    plot_histogram(p3, p3.split(' ', 1)[1], y_label, p3.rsplit(' ', 1)[0], p3_file, x_max)

    print("Histograms have been created and saved as PNG files.")
    return p1_file, p2_file, p3_file

def find_failing_parts(file_path, column_names, Idss_lim):
    df = pd.read_csv(file_path)
    
    def generate_failure_string(die_col, idss_col, limit):
        failures = df[df[idss_col] > limit]
        failure_string = "\n".join(f"DIE ID: {int(row[die_col])}    Idss: {row[idss_col]:.5f}" 
                                   for idx, row in failures.iterrows())
        print(failure_string)
        return failure_string
    
    test_1_failures = generate_failure_string(column_names[0], column_names[1], Idss_lim)
    test_2_failures = generate_failure_string(column_names[2], column_names[3], Idss_lim)
    test_3_failures = generate_failure_string(column_names[4], column_names[5], Idss_lim)
    
    return test_1_failures, test_2_failures, test_3_failures

def create_pptx_report(report_title, report_subtitle, summary_report, column_names,
                       test_1_report, test_2_report, test_3_report, 
                       plot_1_file, plot_2_file, plot_3_file, sub_dir):
    # Create a presentation object
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = report_title
    subtitle.text = report_subtitle

    # Executive Summary Slide
    executive_summary_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(executive_summary_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"

    # Read the summary report text from file
    with open(summary_report, "r") as file:
        summary_text = file.read()

    text_box = slide.shapes.placeholders[1].text_frame
    text_box.text = summary_text
    for paragraph in text_box.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)  #the text size

    # Function to add image and wrapped text to a slide
    def add_image_and_text(slide, img_path, text, left, top, height):
        slide.shapes.add_picture(img_path, left, top, height=height)
        
        text_box = slide.shapes.add_textbox(left - 1.5, Inches(5.2), Inches(8.0), Inches(1.0))
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True  # Enable word wrap
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)  # Set text size
                run.font.color.rgb = RGBColor(0, 0, 0)  # Set text color to black

    # Test #1 Slide
    slide_1_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_1_layout)
    title = slide.shapes.title
    title.text = f"{column_names[1].split(' ', 1)[0]} Test Results"

    # Add image for Test #1
    img_path = plot_1_file
    with open(test_1_report, "r") as file:
        slide_1_text = file.read()
    add_image_and_text(slide, img_path, slide_1_text, Inches(2.5), Inches(1.5), Inches(3.5))

    # Test #2 Slide
    slide_2_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_2_layout)
    title = slide.shapes.title
    title.text = f"{column_names[3].split(' ', 1)[0]} Test Results"

    # Add image for Test #2
    img_path = plot_2_file
    with open(test_2_report, "r") as file:
        slide_2_text = file.read()
    add_image_and_text(slide, img_path, slide_2_text, Inches(2.5), Inches(1.5), Inches(3.5))

    # Test #3 Slide
    slide_3_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_3_layout)
    title = slide.shapes.title
    title.text = f"{column_names[5].split(' ', 1)[0]} Test Results"

    # Add image for Test #3
    img_path = plot_3_file
    with open(test_3_report, "r") as file:
        slide_3_text = file.read()
    add_image_and_text(slide, img_path, slide_3_text, Inches(2.5), Inches(1.5), Inches(3.5))

    # Save the presentation
    pptx_file = path_to_save_file(sub_dir, report_title + '.pptx', False)
    prs.save(pptx_file)

def main():
    print("\n\n## Welcome to the Auto REL Report Generator ##")
    print('-------------------------------')

    sub_dir = 'data'  # Specify the sub-directory to save the file

    Idss_lim = 1  # Set the limit for Idss values

    report_title = "BMXD-456 Semiconductor Summary Report"
    report_subtitle = "The AI Semiconductor Company"

    # Call the create_sample_data function
    file_path, column_names = create_sample_data(sub_dir=sub_dir)  

    # Call the plot_idss_values function
    plot_1_file, plot_2_file, plot_3_file = plot_idss_values(file_path, column_names, 
                                                             sub_dir, Idss_lim) 

    test_1_failures, test_2_failures, test_3_failures =  find_failing_parts(file_path, 
                                                                            column_names, 
                                                                            Idss_lim)
    
    #llm defaults to gpt-4o
    llm = ChatOpenAI(model="gpt-4o", temperature=0.7)

    # Define the paths to save the reports
    test_1_report = path_to_save_file(sub_dir, column_names[1] + '.txt', False)
    test_2_report = path_to_save_file(sub_dir, column_names[3] + '.txt', False)
    test_3_report = path_to_save_file(sub_dir, column_names[5] + '.txt', False)
    summary_report = path_to_save_file(sub_dir, 'summary_report.txt', False)

    # Call the crew function
    inputs = {'test_1': column_names[1], 'test_1_failures': test_1_failures,
              'test_2': column_names[3], 'test_2_failures': test_2_failures,
              'test_3': column_names[5], 'test_3_failures': test_3_failures,
              'Idss_lim': Idss_lim}
    RelReportCrew(test_1_report, test_2_report, test_3_report, summary_report, llm
                  ).crew().kickoff(inputs=inputs)
    
    # Call the create_pptx_report function
    create_pptx_report(report_title, report_subtitle, summary_report, column_names,
                       test_1_report, test_2_report, test_3_report, 
                       plot_1_file, plot_2_file, plot_3_file, sub_dir)

if __name__ == "__main__":
    main()  # Call the main function